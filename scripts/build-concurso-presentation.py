#!/usr/bin/env python3
"""Generate branded Concurso 2026 presentation assets (PPTX + PDF)."""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
RECURSOS = ROOT / "RECURSOS"
ARTIFACT = ROOT / "backend" / "ml" / "artifacts" / "randomforest-outbreak-v1.0.0.json"
PORTADA = RECURSOS / "portada.png"
BG_CONTENT = RECURSOS / ".bg-content.png"
BG_CLOSING = RECURSOS / ".bg-closing.png"

TEAM_ID = "200"
TEAM_LEVEL = "Intermedio (IA)"
REPO_URL = "https://github.com/luisrapalino/concurso-datos-abiertos-colombia-2026"
DEMO_URL = "https://concurso-datos-abiertos-colombia-20.vercel.app/brotes"
API_URL = "https://epintel-api.onrender.com/docs"

# Design tokens — docs/DESIGN.md
C_DEEP = (0x0A, 0x3D, 0x38)
C_PRIMARY = (0x0A, 0x5C, 0x54)
C_ACCENT = (0x5E, 0xEA, 0xD4)
C_TEXT = (0x14, 0x28, 0x24)
C_MUTED = (0x5C, 0x6F, 0x6B)
C_LIGHT = (0xF0, 0xF5, 0xF3)
C_CARD = (0xFA, 0xFC, 0xFB)
C_SIGNAL = (0xB4, 0x53, 0x09)
C_WHITE = (0xFF, 0xFF, 0xFF)
C_BORDER = (0xC8, 0xD9, 0xD4)

SLIDE_W = 13.333
SLIDE_H = 7.5


def _load_metrics() -> dict:
    if ARTIFACT.exists():
        return json.loads(ARTIFACT.read_text(encoding="utf-8"))
    return {}


def _ensure(package: str, pip_name: str | None = None) -> None:
    try:
        __import__(package)
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name or package, "-q"])


def _hex(color: tuple[int, int, int]) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def _lerp(a: tuple[int, int, int], b: tuple[int, int, int], t: float) -> tuple[int, int, int]:
    return tuple(round(a[i] + (b[i] - a[i]) * t) for i in range(3))


def _vgradient(width: int, height: int, top: tuple[int, int, int], bottom: tuple[int, int, int]):
    from PIL import Image

    gradient = Image.new("RGB", (1, height))
    for y in range(height):
        gradient.putpixel((0, y), _lerp(top, bottom, y / max(height - 1, 1)))
    return gradient.resize((width, height))


def _radial_glow(width: int, height: int, center, radius: int, color: tuple[int, int, int], max_alpha: int):
    from PIL import Image, ImageDraw

    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow)
    steps = 60
    for step in range(steps, 0, -1):
        ratio = step / steps
        r = int(radius * ratio)
        alpha = int(max_alpha * (1 - ratio) ** 1.6)
        gdraw.ellipse(
            (center[0] - r, center[1] - r, center[0] + r, center[1] + r),
            fill=(*color, alpha),
        )
    return glow


def _draw_radar(layer, center, max_radius: int, *, ring_alpha: int = 60, sweep: bool = True, signals: bool = True) -> None:
    from PIL import ImageDraw
    import math

    draw = ImageDraw.Draw(layer, "RGBA")
    rings = 7
    for index in range(1, rings + 1):
        r = int(max_radius * index / rings)
        alpha = int(ring_alpha * (index / rings))
        draw.ellipse(
            (center[0] - r, center[1] - r, center[0] + r, center[1] + r),
            outline=(*C_ACCENT, alpha),
            width=2,
        )
    for angle_deg in range(0, 360, 30):
        angle = math.radians(angle_deg)
        x2 = center[0] + max_radius * math.cos(angle)
        y2 = center[1] + max_radius * math.sin(angle)
        draw.line((center[0], center[1], x2, y2), fill=(*C_ACCENT, 22), width=1)
    if sweep:
        for step in range(70):
            angle = math.radians(-35 - step * 0.9)
            x2 = center[0] + max_radius * math.cos(angle)
            y2 = center[1] + max_radius * math.sin(angle)
            alpha = int(70 * (1 - step / 70))
            draw.line((center[0], center[1], x2, y2), fill=(*C_ACCENT, alpha), width=2)
    if not signals:
        return
    signal_dots = [
        (0.55, -50, C_ACCENT, 16),
        (0.78, 120, C_SIGNAL, 20),
        (0.34, 200, C_ACCENT, 12),
        (0.9, -140, C_SIGNAL, 14),
    ]
    for frac, ang, color, size in signal_dots:
        angle = math.radians(ang)
        r = max_radius * frac
        cx = center[0] + r * math.cos(angle)
        cy = center[1] + r * math.sin(angle)
        draw.ellipse((cx - size, cy - size, cx + size, cy + size), fill=(*color, 70))
        draw.ellipse((cx - size / 2, cy - size / 2, cx + size / 2, cy + size / 2), fill=(*color, 235))


def _font_path(name: str) -> str | None:
    candidates = [
        f"/usr/share/fonts/truetype/dejavu/{name}",
        f"/usr/local/lib/python3.12/site-packages/matplotlib/mpl-data/fonts/ttf/{name}",
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    try:
        import matplotlib

        mpl_dir = Path(matplotlib.get_data_path()) / "fonts" / "ttf" / name
        if mpl_dir.exists():
            return str(mpl_dir)
    except Exception:
        pass
    return None


def _font(bold: bool, size: int):
    from PIL import ImageFont

    name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    path = _font_path(name) or _font_path("DejaVuSans.ttf")
    if path:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            pass
    return ImageFont.load_default()


def _spaced(text: str, spacing: str = " ") -> str:
    return spacing.join(list(text))


def _generate_visual_assets() -> None:
    """Create branded backgrounds: hero portada, content and closing surfaces."""
    from PIL import Image, ImageDraw

    RECURSOS.mkdir(parents=True, exist_ok=True)
    width, height = 1920, 1080

    # ---- Portada (hero) ----
    base = _vgradient(width, height, (0x05, 0x24, 0x20), (0x0A, 0x5C, 0x54)).convert("RGBA")
    base = Image.alpha_composite(base, _radial_glow(width, height, (1430, 520), 720, C_PRIMARY, 120))
    radar = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    _draw_radar(radar, (1440, 540), 470)
    base = Image.alpha_composite(base, radar)

    draw = ImageDraw.Draw(base, "RGBA")
    draw.rectangle((0, 0, width, 12), fill=(*C_ACCENT, 255))
    draw.rectangle((0, 0, 14, height), fill=(*C_ACCENT, 90))

    eyebrow_font = _font(True, 30)
    title_font = _font(True, 132)
    sub_font = _font(False, 40)
    pill_font = _font(True, 27)
    contact_font = _font(False, 27)

    draw.text((110, 250), _spaced("CONCURSO DATOS AL ECOSISTEMA 2026"), fill=(*C_ACCENT, 255), font=eyebrow_font)
    draw.text((104, 315), "Radar de Brotes", fill=(*C_WHITE, 255), font=title_font)
    draw.rectangle((112, 500, 360, 512), fill=(*C_ACCENT, 255))
    draw.text((110, 545), "Inteligencia epidemiológica territorial", fill=(*C_WHITE, 235), font=sub_font)
    draw.text((110, 600), "con datos abiertos e IA explicable", fill=(0xB8, 0xD8, 0xD1, 255), font=sub_font)

    pills = [f"Equipo ID {TEAM_ID}", TEAM_LEVEL, "Salud y Bienestar"]
    x = 110
    for pill in pills:
        bbox = draw.textbbox((0, 0), pill, font=pill_font)
        w_text = bbox[2] - bbox[0]
        pill_w = w_text + 56
        draw.rounded_rectangle((x, 720, x + pill_w, 772), radius=26, outline=(*C_ACCENT, 220), width=2)
        draw.text((x + 28, 732), pill, fill=(*C_ACCENT, 255), font=pill_font)
        x += pill_w + 22

    draw.text((110, 940), "luisrapalino88@gmail.com   ·   Rapalinorokate@gmail.com", fill=(0x9F, 0xC8, 0xC0, 255), font=contact_font)
    draw.text((110, 985), REPO_URL, fill=(0x7C, 0xB0, 0xA7, 255), font=contact_font)
    base.convert("RGB").save(PORTADA, "PNG", optimize=True)

    # ---- Content background ----
    content = Image.new("RGBA", (width, height), (*C_LIGHT, 255))
    band = _vgradient(width, 168, C_DEEP, C_PRIMARY).convert("RGBA")
    content.paste(band, (0, 0))
    cdraw = ImageDraw.Draw(content, "RGBA")
    cdraw.rectangle((0, 168, width, 176), fill=(*C_ACCENT, 255))
    cdraw.rectangle((0, 0, 14, height), fill=(*C_PRIMARY, 255))
    faint = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    _draw_radar(faint, (1780, 1060), 560, ring_alpha=16, sweep=False, signals=False)
    content = Image.alpha_composite(content, faint)
    content.convert("RGB").save(BG_CONTENT, "PNG", optimize=True)

    # ---- Closing background ----
    closing = _vgradient(width, height, (0x06, 0x26, 0x22), (0x0A, 0x4A, 0x44)).convert("RGBA")
    closing = Image.alpha_composite(closing, _radial_glow(width, height, (960, 620), 760, C_PRIMARY, 90))
    cradar = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    _draw_radar(cradar, (960, 560), 620, ring_alpha=32)
    closing = Image.alpha_composite(closing, cradar)
    cldraw = ImageDraw.Draw(closing, "RGBA")
    cldraw.rectangle((0, 0, width, 12), fill=(*C_ACCENT, 255))
    closing.convert("RGB").save(BG_CLOSING, "PNG", optimize=True)


def _rgb(color: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*color)


def _set_bg(slide, color: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, left, top, width, height, color: tuple[int, int, int], *, line: bool = False):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if not line:
        shape.line.fill.background()
    return shape


def _rounded_card(slide, left, top, width, height, accent: tuple[int, int, int]) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Pt

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(C_CARD)
    card.line.color.rgb = _rgb(C_BORDER)
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    _rect(slide, left, top + Emu(int(height * 0.16)), Emu(int(width * 0.028)), Emu(int(height * 0.68)), accent)


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


_SLIDE_SEQ = {"n": 0}


def _add_background(slide, image: Path) -> None:
    from pptx.util import Inches

    _set_bg(slide, C_LIGHT)
    if image.exists():
        pic = slide.shapes.add_picture(str(image), Inches(0), Inches(0), width=Inches(SLIDE_W), height=Inches(SLIDE_H))
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)


def _add_brand_header(slide, title: str, *, eyebrow: str | None = None) -> None:
    from pptx.util import Inches, Pt

    _add_background(slide, BG_CONTENT)
    _SLIDE_SEQ["n"] += 1

    if eyebrow:
        tbe = _textbox(slide, Inches(0.6), Inches(0.16), Inches(9), Inches(0.35))
        tfe = tbe.text_frame
        tfe.text = eyebrow.upper()
        run = tfe.paragraphs[0].runs[0]
        run.font.size = Pt(11)
        run.font.color.rgb = _rgb(C_ACCENT)
        run.font.bold = True
        top = 0.5
    else:
        top = 0.28

    tb = _textbox(slide, Inches(0.6), Inches(top), Inches(10.5), Inches(0.75))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = _rgb(C_WHITE)
    p.font.name = "Calibri"

    tb2 = _textbox(slide, Inches(11.4), Inches(0.34), Inches(1.6), Inches(0.6))
    tf2 = tb2.text_frame
    tf2.text = f"{_SLIDE_SEQ['n']:02d}"
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = _rgb(C_ACCENT)
    tf2.paragraphs[0].alignment = 2


def _add_footer(slide, text: str = "Radar de Brotes · Datos abiertos + IA explicable") -> None:
    from pptx.util import Inches, Pt

    tb = _textbox(slide, Inches(0.6), Inches(7.08), Inches(9), Inches(0.35))
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = _rgb(C_MUTED)

    tb2 = _textbox(slide, Inches(10.3), Inches(7.08), Inches(2.7), Inches(0.35))
    tf2 = tb2.text_frame
    tf2.text = f"Equipo ID {TEAM_ID}"
    tf2.paragraphs[0].font.size = Pt(9)
    tf2.paragraphs[0].font.color.rgb = _rgb(C_MUTED)
    tf2.paragraphs[0].alignment = 2


def _add_bullets(slide, left, top, width, height, bullets: list[str], *, size: int = 16) -> None:
    from pptx.util import Pt

    tb = _textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"■  {bullet}"
        p.space_after = Pt(10)
        p.font.size = Pt(size)
        p.font.color.rgb = _rgb(C_TEXT)
        p.font.name = "Calibri"
        if p.runs:
            p.runs[0].font.color.rgb = _rgb(C_ACCENT)
            if len(p.runs) > 1:
                p.runs[1].font.color.rgb = _rgb(C_TEXT)


def _slide_title_hero(prs) -> None:
    from pptx.util import Inches, Pt

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, C_DEEP)

    if PORTADA.exists():
        slide.shapes.add_picture(str(PORTADA), Inches(0), Inches(0), width=Inches(SLIDE_W), height=Inches(SLIDE_H))
    else:
        _rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.12), C_ACCENT)
        tb = _textbox(slide, Inches(0.85), Inches(1.6), Inches(8.5), Inches(1.2))
        tf = tb.text_frame
        tf.text = "Radar de Brotes"
        tf.paragraphs[0].font.size = Pt(54)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _rgb(C_WHITE)


def _slide_content(prs, title: str, bullets: list[str], *, highlight: str | None = None, eyebrow: str | None = None) -> None:
    from pptx.util import Inches, Pt

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_brand_header(slide, title, eyebrow=eyebrow)

    if highlight:
        _rounded_card(slide, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.9), C_ACCENT)
        tb = _textbox(slide, Inches(0.95), Inches(2.16), Inches(11.4), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = highlight
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = _rgb(C_PRIMARY)
        tf.paragraphs[0].font.bold = True
        start_y = 3.2
    else:
        start_y = 2.15

    _add_bullets(
        slide,
        Inches(0.9),
        Inches(start_y),
        Inches(11.5),
        Inches(4.4),
        bullets,
        size=17,
    )
    _add_footer(slide)


def _slide_data_sources(prs) -> None:
    from pptx.util import Inches, Pt

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_brand_header(slide, "Datos abiertos", eyebrow="Fuentes oficiales")

    sources = [
        ("SIVIGILA", "Casos semanales: dengue, hepatitis, chikungunya, fiebre tifoidea"),
        ("Vacunación", "Cobertura pentavalente DPT-HepB-Hib (datos.gov.co)"),
        ("Calidad del aire", "PM2.5 anual municipal con fallback SISAIRE"),
        ("Indicadores INS", "Mortalidad general y partos institucionales"),
        ("DIVIPOLA", "Validación territorial embebida"),
        ("Panel ML", "20 ciudades · ~9.000 obs. · 15 variables · PostgreSQL"),
    ]
    positions = [(0.6, 1.6), (4.4, 1.6), (8.2, 1.6), (0.6, 3.75), (4.4, 3.75), (8.2, 3.75)]
    colors = [C_PRIMARY, C_SIGNAL, C_PRIMARY, C_SIGNAL, C_PRIMARY, C_ACCENT]

    for (x, y), (title, desc), accent in zip(positions, sources, colors, strict=True):
        _rounded_card(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.85), accent)
        tb = _textbox(slide, Inches(x + 0.18), Inches(y + 0.28), Inches(3.15), Inches(0.45))
        tf = tb.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _rgb(C_PRIMARY)
        tb2 = _textbox(slide, Inches(x + 0.18), Inches(y + 0.72), Inches(3.15), Inches(0.95))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.text = desc
        tf2.paragraphs[0].font.size = Pt(11)
        tf2.paragraphs[0].font.color.rgb = _rgb(C_MUTED)

    _add_footer(slide)


def _slide_architecture(prs) -> None:
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches, Pt

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_brand_header(slide, "Solución e inteligencia artificial", eyebrow="Arquitectura")

    boxes = [
        (0.75, 1.7, "Datos abiertos", "Socrata · INS · SIVIGILA"),
        (3.55, 1.7, "Ingestión", "Worker + PostgreSQL"),
        (6.35, 1.7, "Modelo ML", "Random Forest + SHAP"),
        (9.15, 1.7, "Plataforma", "FastAPI + Next.js"),
    ]
    for x, y, title, sub in boxes:
        _rounded_card(slide, Inches(x), Inches(y), Inches(2.45), Inches(1.35), C_PRIMARY)
        tb = _textbox(slide, Inches(x + 0.15), Inches(y + 0.35), Inches(2.15), Inches(0.5))
        tf = tb.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _rgb(C_PRIMARY)
        tb2 = _textbox(slide, Inches(x + 0.15), Inches(y + 0.78), Inches(2.15), Inches(0.45))
        tb2.text_frame.text = sub
        tb2.text_frame.paragraphs[0].font.size = Pt(10)
        tb2.text_frame.paragraphs[0].font.color.rgb = _rgb(C_MUTED)

    for index in range(3):
        x1 = Inches(0.75 + 2.45 + index * 2.8)
        y = Inches(2.35)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x1 + Inches(0.35), y)
        conn.line.color.rgb = _rgb(C_ACCENT)
        conn.line.width = Pt(2.5)

    bullets = [
        "Unidad analítica: municipio × semana epidemiológica × dengue (código 210)",
        "Validación temporal: entrenamiento ≤ 2020 · prueba ≥ 2021",
        "Explicabilidad SHAP TreeExplainer · fallback rule-based auditable",
        "Despliegue reproducible con Docker Compose",
    ]
    _add_bullets(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(3.1), bullets, size=15)
    _add_footer(slide)


def _slide_metrics(prs, metrics: dict) -> None:
    from pptx.util import Inches, Pt

    validation = metrics.get("validation", {})
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_brand_header(slide, "Resultados y demo", eyebrow="Evidencia")

    cards = [
        (str(metrics.get("training_samples", "—")), "Muestras ML", C_PRIMARY),
        (f"{validation.get('train_samples', '—')} / {validation.get('test_samples', '—')}", "Train / Test", C_SIGNAL),
        (str(validation.get("f1", "—")), "F1 temporal", C_PRIMARY),
    ]
    x = 0.75
    for value, label, color in cards:
        _rounded_card(slide, Inches(x), Inches(1.6), Inches(3.7), Inches(1.6), color)
        tb = _textbox(slide, Inches(x + 0.35), Inches(1.78), Inches(3.1), Inches(0.9))
        tf = tb.text_frame
        tf.text = value
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _rgb(color)
        tb2 = _textbox(slide, Inches(x + 0.35), Inches(2.62), Inches(3.1), Inches(0.4))
        tb2.text_frame.text = label
        tb2.text_frame.paragraphs[0].font.size = Pt(12)
        tb2.text_frame.paragraphs[0].font.color.rgb = _rgb(C_MUTED)
        x += 4.0

    bullets = [
        "Modelo promovido: randomforest-outbreak-v1.0.0 (Random Forest + SHAP)",
        "Cobertura: 20 ciudades principales · ~9.000 observaciones curadas",
        "Cautela: panel acotado; validación epidemiológica externa pendiente",
        f"Demo en vivo: {DEMO_URL}  ·  Mapa /informe  ·  API {API_URL}",
    ]
    _add_bullets(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(3.0), bullets, size=15)
    _add_footer(slide)


def _slide_closing(prs) -> None:
    from pptx.util import Inches, Pt

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, C_DEEP)
    _add_background(slide, BG_CLOSING)

    tbe = _textbox(slide, Inches(0.9), Inches(1.0), Inches(9), Inches(0.4))
    tbe.text_frame.text = _spaced("CIERRE").upper()
    tbe.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    tbe.text_frame.paragraphs[0].runs[0].font.bold = True
    tbe.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_ACCENT)

    tb = _textbox(slide, Inches(0.9), Inches(1.55), Inches(11.2), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Priorizamos municipios con señal de brote usando datos abiertos "
        "de Colombia e inteligencia artificial explicable."
    )
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _rgb(C_WHITE)

    tb2 = _textbox(slide, Inches(0.9), Inches(3.5), Inches(11), Inches(2.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    lines = [
        "Integración multifuente · Explicabilidad SHAP · Plataforma demostrable",
        REPO_URL,
        "Gracias — preguntas técnicas bienvenidas",
    ]
    for index, line in enumerate(lines):
        p = tf2.paragraphs[0] if index == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if index else 18)
        p.font.color.rgb = _rgb(C_ACCENT if index == 2 else C_WHITE)
        p.space_after = Pt(12)


def build_pptx(metrics: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _slide_title_hero(prs)
    _slide_content(
        prs,
        "Problema y objetivo",
        [
            "Los brotes de enfermedades transmisibles exigen detección temprana, pero la información "
            "está dispersa en fuentes abiertas difíciles de cruzar en tiempo útil.",
            "Objetivo: priorizar municipios con señal de brote integrando morbilidad, vacunación, "
            "calidad del aire y acceso a salud.",
            "Enfoque: apoyo analítico con revisión humana — no activación automática de emergencias.",
        ],
        highlight="¿Cómo anticipar un brote antes de que el sistema de salud se sobrecargue?",
        eyebrow="Contexto",
    )
    _slide_data_sources(prs)
    _slide_architecture(prs)
    _slide_metrics(prs, metrics)
    _slide_content(
        prs,
        "Impacto y sostenibilidad",
        [
            "Beneficiarios: equipos de vigilancia en salud pública y gestión territorial",
            "Permite focalizar recursos antes del pico de casos",
            "Explicación clara de factores — evita caja negra",
            "Reproducible con datos abiertos oficiales de Colombia",
            "Sostenibilidad: ingestión automática, documentación, CI y despliegue Docker",
        ],
        highlight="Valor público: prevención y respuesta temprana con trazabilidad",
        eyebrow="Impacto",
    )
    _slide_content(
        prs,
        "Repositorio y recursos",
        [
            f"GitHub: {REPO_URL}",
            "README.md — ficha técnica y guía de reproducción",
            "docs/concurso-alignment.md — trazabilidad del reto",
            "docs/data-dictionary.md — 15 variables del modelo",
            "docs/ml-evaluation.md — validación, métricas y límites",
            "RECURSOS/ — presentación, PDF y portada",
        ],
        eyebrow="Entregables",
    )
    _slide_closing(prs)

    out = RECURSOS / "Presentacion.pptx"
    prs.save(out)
    return out


def build_pdf(metrics: dict) -> Path:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        Image,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    validation = metrics.get("validation", {})
    out = RECURSOS / "presentacion.pdf"
    page_w, page_h = landscape(A4)

    def _draw_brand(canvas, doc, *, cover: bool = False) -> None:
        canvas.saveState()
        if cover:
            canvas.setFillColor(colors.HexColor("#0a3d38"))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
            canvas.setFillColor(colors.HexColor("#5eead4"))
            canvas.rect(0, page_h - 0.35 * cm, page_w, 0.35 * cm, fill=1, stroke=0)
        else:
            canvas.setFillColor(colors.HexColor("#f0f5f3"))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
            canvas.setFillColor(colors.HexColor("#0a5c54"))
            canvas.rect(0, page_h - 1.6 * cm, page_w, 1.6 * cm, fill=1, stroke=0)
            canvas.setFillColor(colors.HexColor("#0a3d38"))
            canvas.rect(0, page_h - 1.6 * cm, 0.45 * cm, 1.6 * cm, fill=1, stroke=0)
            canvas.setFillColor(colors.HexColor("#5eead4"))
            canvas.rect(0, page_h - 0.25 * cm, page_w, 0.08 * cm, fill=1, stroke=0)
            canvas.setFillColor(colors.HexColor("#5c6f6b"))
            canvas.setFont("Helvetica", 8)
            canvas.drawString(1.8 * cm, 0.55 * cm, "Radar de Brotes · Datos abiertos + IA explicable")
            canvas.setFillColor(colors.HexColor("#5eead4"))
            canvas.setFont("Helvetica-Bold", 9)
            canvas.drawRightString(page_w - 1.8 * cm, 0.55 * cm, f"Equipo {TEAM_ID}")
        canvas.restoreState()

    doc = SimpleDocTemplate(
        str(out),
        pagesize=landscape(A4),
        leftMargin=1.8 * cm,
        rightMargin=1.8 * cm,
        topMargin=2.0 * cm,
        bottomMargin=1.4 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "BrandTitle",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=20,
        textColor=colors.HexColor("#ffffff"),
        spaceAfter=8,
    )
    slide_title = ParagraphStyle(
        "SlideTitle",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        textColor=colors.HexColor("#0a3d38"),
        spaceAfter=10,
    )
    body_style = ParagraphStyle(
        "BrandBody",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=11,
        leading=15,
        textColor=colors.HexColor("#142824"),
        spaceAfter=6,
        leftIndent=12,
        bulletIndent=0,
    )
    accent_style = ParagraphStyle(
        "Accent",
        parent=body_style,
        textColor=colors.HexColor("#0a5c54"),
        fontName="Helvetica-Bold",
        leftIndent=0,
    )
    quote_style = ParagraphStyle(
        "Quote",
        parent=accent_style,
        backColor=colors.HexColor("#fafcfb"),
        borderPadding=8,
        borderColor=colors.HexColor("#5eead4"),
        borderWidth=2,
        leftIndent=8,
    )

    def bullets(lines: list[str]) -> list:
        flow = []
        for line in lines:
            flow.append(Paragraph(f"<font color='#0a5c54'>&#8226;</font>&nbsp;&nbsp;{line}", body_style))
        return flow

    story: list = []

    if PORTADA.exists():
        cover_img = Image(str(PORTADA), width=page_w - 3.6 * cm, height=(page_w - 3.6 * cm) * 9 / 16)
        cover_img.hAlign = "CENTER"
        story.append(Spacer(1, 0.4 * cm))
        story.append(cover_img)
    else:
        story.append(Paragraph("Radar de Brotes", title_style))
        story.append(Paragraph(f"Equipo ID {TEAM_ID} · {TEAM_LEVEL}", accent_style))

    story.append(PageBreak())
    story.append(Paragraph("Problema y objetivo", slide_title))
    story.append(Paragraph("¿Cómo anticipar un brote antes de que el sistema de salud se sobrecargue?", quote_style))
    story.append(Spacer(1, 0.2 * cm))
    story += bullets(
        [
            "Detección temprana con datos dispersos en fuentes abiertas.",
            "Priorizar municipios integrando morbilidad, vacunación, aire y acceso a salud.",
            "Apoyo analítico con revisión humana — no activación automática de emergencias.",
        ]
    )
    story.append(PageBreak())

    story.append(Paragraph("Datos abiertos", slide_title))
    data_table = Table(
        [
            ["SIVIGILA", "Casos semanales (dengue, hepatitis, chikungunya…)"],
            ["Vacunación", "Pentavalente DPT-HepB-Hib"],
            ["PM2.5", "Calidad del aire municipal"],
            ["INS", "Mortalidad y partos institucionales"],
            ["Panel", "20 ciudades · ~9.000 obs. · 15 variables"],
        ],
        colWidths=[4 * cm, 14 * cm],
    )
    data_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#0a5c54")),
                ("TEXTCOLOR", (0, 0), (0, -1), colors.white),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("BACKGROUND", (1, 0), (1, -1), colors.HexColor("#fafcfb")),
                ("TEXTCOLOR", (1, 0), (1, -1), colors.HexColor("#142824")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#c8d9d4")),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    story.append(data_table)
    story.append(PageBreak())

    story.append(Paragraph("Solución e inteligencia artificial", slide_title))
    story += bullets(
        [
            "FastAPI (DDD) + Next.js + Docker Compose",
            "Random Forest + SHAP TreeExplainer",
            "Validación temporal: entrenamiento ≤ 2020 · prueba ≥ 2021",
            "Fallback rule-based auditable si no hay modelo promovido",
        ]
    )
    story.append(PageBreak())

    story.append(Paragraph("Resultados y demo", slide_title))
    metric_data = [
        ["Muestras ML", str(metrics.get("training_samples", "—"))],
        ["Train / Test", f"{validation.get('train_samples', '-')} / {validation.get('test_samples', '-')}"],
        ["F1 temporal", str(validation.get("f1", "—"))],
    ]
    metric_table = Table(metric_data, colWidths=[5.5 * cm, 5.5 * cm], hAlign="LEFT")
    metric_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0a5c54")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#fafcfb")),
                ("TEXTCOLOR", (0, 1), (-1, -1), colors.HexColor("#142824")),
                ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#c8d9d4")),
                ("FONTSIZE", (0, 0), (-1, -1), 11),
                ("TOPPADDING", (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ]
        )
    )
    story.append(metric_table)
    story.append(Spacer(1, 0.25 * cm))
    story.append(Paragraph(f"Demo: {DEMO_URL}", accent_style))
    story.append(PageBreak())

    story += [Paragraph("Impacto y sostenibilidad", slide_title)]
    story += bullets(
        [
            "Equipos de vigilancia en salud pública y gestión territorial",
            "Focalización de recursos antes del pico de casos",
            "Explicabilidad — evita caja negra",
            "Reproducible con datos.gov.co · ingestión automática y CI",
        ]
    )
    story.append(PageBreak())

    story += [Paragraph("Repositorio", slide_title)]
    story += bullets(
        [
            REPO_URL,
            "docs/concurso-alignment.md · docs/data-dictionary.md · docs/ml-evaluation.md",
        ]
    )
    story.append(PageBreak())

    story.append(Paragraph("Cierre", slide_title))
    story.append(
        Paragraph(
            "Priorizamos municipios con señal de brote usando datos abiertos de Colombia "
            "e inteligencia artificial explicable.",
            accent_style,
        )
    )
    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph("Gracias — preguntas técnicas bienvenidas", body_style))

    def on_page(canvas, doc_obj):
        _draw_brand(canvas, doc_obj, cover=False)

    def on_cover(canvas, doc_obj):
        _draw_brand(canvas, doc_obj, cover=True)

    doc.build(story, onFirstPage=on_cover, onLaterPages=on_page)
    return out


def main() -> int:
    RECURSOS.mkdir(parents=True, exist_ok=True)
    metrics = _load_metrics()
    _ensure("PIL", "Pillow")
    _ensure("pptx", "python-pptx")
    _ensure("reportlab")
    _generate_visual_assets()
    pptx_path = build_pptx(metrics)
    pdf_path = build_pdf(metrics)
    for temp in (BG_CONTENT, BG_CLOSING):
        if temp.exists():
            temp.unlink()
    print(f"Generated: {pptx_path}")
    print(f"Generated: {pdf_path}")
    print(f"Generated: {PORTADA}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
